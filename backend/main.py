import asyncio
import re
from fastapi import FastAPI, HTTPException
from fastapi.responses import FileResponse
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel, Field
from typing import List, Literal

from langchain_aws import ChatBedrock
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import JsonOutputParser

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
import logging
import sys

# --- 1. アプリケーションとロギングの初期化 ---
logging.basicConfig(stream=sys.stdout, level=logging.INFO)
app = FastAPI()
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"], # すべてのオリジンを許可する
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# --- 2. データモデル定義 ---
class SlidePlan(BaseModel):
    topic: str = Field(description="このスライドで説明すべき具体的なトピック")
    slide_type: Literal["text_slide", "table_slide"] = Field(description="トピックの内容に適したスライドのタイプ")

class PresentationPlan(BaseModel):
    plan: List[SlidePlan]
    rationale: str = Field(description="生成された構成案全体の論理的な妥当性やストーリーの流れに関する解説")

class TextSlideContent(BaseModel):
    title: str = Field(description="スライドのタイトル")
    content: str = Field(description="箇条書き形式のスライド本文")

class TableData(BaseModel):
    headers: List[str] = Field(description="表のヘッダーリスト")
    rows: List[List[str]] = Field(description="表の各行のデータリスト（リストのリスト）")

class TableSlideContent(BaseModel):
    title: str = Field(description="スライドのタイトル")
    table_data: TableData = Field(description="スライドに表示する表データ")

class UserInput(BaseModel):
    prompt: str

# --- 3. LLMと専門エージェント関数の定義 ---
llm = ChatBedrock(
#    model_id="us.anthropic.claude-sonnet-4-20250514-v1:0",
    model_id="us.anthropic.claude-3-7-sonnet-20250219-v1:0",
#    model_id="us.anthropic.claude-3-5-sonnet-20241022-v2:0",
#    model_id="us.anthropic.claude-3-5-sonnet-20240620-v1:0",
#    model_id="us.amazon.nova-pro-v1:0",
#    model_id="us.amazon.nova-premier-v1:0",
#    model_id="us.meta.llama3-3-70b-instruct-v1:0",
    model_kwargs={"temperature": 0.1},
)


async def run_supervisor_agent(user_prompt: str) -> PresentationPlan:
    """プレゼン全体の計画を立案するエージェント関数"""
    parser = JsonOutputParser(pydantic_object=PresentationPlan)
    prompt = ChatPromptTemplate.from_messages([
        ("system", """あなたはプレゼン構成のプロフェッショナルです。ユーザーの要望に基づき、最高のプレゼンテーション構成を提案してください。

# 指示
1.  必ず「序論・本論・結論」のフレームワークに沿って構成を考えてください。
    - **序論**: プレゼンの目的、背景、問題提起、聞き手がこれから何を得られるかを簡潔に述べます。
    - **本論**: 主張や提案の核心部分です。データ、事実、比較を用いて論理的に説明します。
    - **結論**: 全体の要約、キーメッセージの再強調、そして聞き手への行動喚起で締めくくります。
2.  各スライドで語るべき`topic`を決定した後、以下の「スライドタイプの選択基準」に従って、最適な`slide_type`を割り当ててください。
3.  最終的なアウトプットとして、構成案（plan）だけでなく、なぜその構成が論理的で説得力があるのかを説明する「根拠（rationale）」も生成してください。

# スライドタイプの選択基準
- `table_slide` は、複数の項目を共通の軸で「比較・対照」する場合に限定して使用してください。例えば、「製品A/B/Cの機能比較」「プラン松/竹/梅の料金比較」「新旧手法のメリット・デメリット対照」などが該当します。
- それ以外の、概念の説明、背景、ストーリー、結論の要約など、情報を順番に説明する場合は `text_slide` を使用してください。

{format_instructions}"""),
        ("human", "{user_prompt}")
    ])
    chain = prompt | llm | parser
    result_dict = await chain.ainvoke({"user_prompt": user_prompt, "format_instructions": parser.get_format_instructions()})
    return PresentationPlan(**result_dict)

async def run_text_agent(topic: str) -> TextSlideContent:
    """テキストスライドの内容を生成するエージェント関数"""
    parser = JsonOutputParser(pydantic_object=TextSlideContent)
    prompt = ChatPromptTemplate.from_messages([
        ("system", "あなたはプレゼンの本文作成の専門家です。与えられたトピックに基づき、スライドのタイトルと本文を作成してください。本文は改行で区切られたプレーンなテキストで生成し、文頭に「・」や「-」などの箇条書き記号は含めないでください。各項目は要点を押さえた簡潔な記述を心がけてください。\n{format_instructions}"),
        ("human", "トピック: {topic}")
    ])
    chain = prompt | llm | parser
    result_dict = await chain.ainvoke({"topic": topic, "format_instructions": parser.get_format_instructions()})
    return TextSlideContent(**result_dict)

async def run_table_agent(topic: str) -> TableSlideContent:
    """表スライドの内容を生成するエージェント関数"""
    parser = JsonOutputParser(pydantic_object=TableSlideContent)
    prompt = ChatPromptTemplate.from_messages([
        ("system", "あなたは情報を表にまとめる専門家です。与えられたトピックに基づき、スライドのタイトルと、内容をまとめた表データを作成してください。\n{format_instructions}"),
        ("human", "トピック: {topic}")
    ])
    chain = prompt | llm | parser
    result_dict = await chain.ainvoke({"topic": topic, "format_instructions": parser.get_format_instructions()})
    return TableSlideContent(**result_dict)


# --- 4. APIエンドポイント定義 ---

@app.post("/api/generate-plan", response_model=PresentationPlan)
async def generate_plan(user_input: UserInput):
    """【計画フェーズ】ユーザーの指示からスライドの構成案（計画）を生成する"""
    print("---API: /api/generate-plan---")
    try:
        plan = await run_supervisor_agent(user_input.prompt)
        return plan
    except Exception as e:
        print(f"Error in /api/generate-plan: {e}")
        raise HTTPException(status_code=500, detail="スライド計画の生成に失敗しました。")


@app.post("/api/create-slides")
async def create_slides(presentation_plan: PresentationPlan):
    """【実行フェーズ】承認された計画に基づき、スライド内容を生成してPPTXファイルを作成する"""
    print("---API: /api/create-slides---")
    try:
        prs = Presentation("template.pptx")
        title_and_content_layout = prs.slide_layouts[4]

        # 処理を単一のループに統合
        for i, slide_plan in enumerate(presentation_plan.plan):
            try:
                # 1. スライド内容を1枚ずつ生成
                if slide_plan.slide_type == "text_slide":
                    result = await run_text_agent(slide_plan.topic)
                elif slide_plan.slide_type == "table_slide":
                    result = await run_table_agent(slide_plan.topic)
                else:
                    continue # 未知のスライドタイプの場合はスキップ
                
                
                # 2. 生成された内容をすぐにPowerPointに書き込む
                slide = prs.slides.add_slide(title_and_content_layout)
                title_shape = slide.shapes.title
                
                if isinstance(result, TextSlideContent):
                    title_shape.text = sanitize_text(result.title) 
                    content_shape = slide.placeholders[1]
                    content_shape.text = sanitize_text(result.content)
                    
                    # テキストに合わせてフォントサイズを自動調整
                    content_shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

                    # フォントを統一
                    set_font_for_shape(title_shape, "BIZ UDPGothic")
                    set_font_for_shape(content_shape, "BIZ UDPGothic")

                elif isinstance(result, TableSlideContent):
                    title_shape.text = sanitize_text(result.title)
                    set_font_for_shape(title_shape, "BIZ UDPGothic")
                    slide.placeholders[1].text_frame.clear() # 本文プレースホルダーをクリア
                    draw_table_on_slide(slide, prs, result.table_data)

                # for debug
                print(f"--- Writing Slide {i+1} ---")
                print(f"Type: {slide_plan.slide_type}")
                if isinstance(result, TextSlideContent):
                    print(f"Title: {repr(result.title)}")
                    print(f"Content: {repr(result.content)}")
                elif isinstance(result, TableSlideContent):
                    print(f"Title: {repr(result.title)}")
                    print(f"Table Data: {result.table_data}")

            except Exception as e:
                raise HTTPException(status_code=500, detail=f"スライド {i + 1} (トピック: '{slide_plan.topic}') の生成に失敗しました: {e}")

            await asyncio.sleep(1) # 1秒待機して次のAPIコールへ

        # 3. ループ完了後、ファイルを保存して返す
        output_filename = "generated_presentation.pptx"
        prs.save(output_filename)
        return FileResponse(output_filename, media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation', filename=output_filename)

    except HTTPException as http_exc:
        raise http_exc
    except Exception as e:
        print(f"Error in /api/create-slides: {e}")
        raise HTTPException(status_code=500, detail=f"PowerPointファイルの生成中に予期せぬエラーが発生しました: {e}")

# --- 5. PowerPoint生成ヘルパー関数 ---

def sanitize_text(text):
    """テキストをサニタイズする最終版：制御文字、特殊ダッシュ、前後の空白を除去"""
    if text is None:
        return ""
    
    clean_text = str(text)

    # 1. 制御文字を削除（既存の処理）
    clean_text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', clean_text)
    
    # 2. 様々な種類のダッシュやハイフンを、標準的な半角ハイフンマイナスに統一
    # U+2010 (HYPHEN) から U+2015 (HORIZONTAL BAR) までの範囲を対象
    dash_chars = r'[\u2010-\u2015]'
    clean_text = re.sub(dash_chars, '-', clean_text)
    
    # 3. 前後の空白を削除（既存の処理）
    return clean_text.strip()


def set_font_for_shape(shape, font_name):
    """指定されたシェイプ内のすべてのテキストにフォントを適用する"""
    if not hasattr(shape, 'text_frame'):
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name

def draw_table_on_slide(slide, prs, table_data: TableData):
    """【改善版】スライドにテーブルを描画する関数"""
    rows_count = len(table_data.rows) + 1
    cols_count = len(table_data.headers)
    
    # 列がない場合はテーブルを作成しない
    if cols_count == 0:
        return

    table_width = Inches(9.0)
    left = (prs.slide_width - table_width) / 2
    top = Inches(1.8)
    
    # 先にダミーの高さでテーブルを作成
    shape = slide.shapes.add_table(rows_count, cols_count, left, top, table_width, Inches(1))
    table = shape.table

    # --- ヘッダーの入力 ---
    for i, header in enumerate(table_data.headers):
        cell = table.cell(0, i)
        cell.text = sanitize_text(header)
        # フォント設定
        p = cell.text_frame.paragraphs[0]
        p.font.name = "BIZ UDPGothic"
        p.font.bold = True
        p.font.size = Pt(14)

    # --- データ行の入力 ---
    for r_idx, row_data in enumerate(table_data.rows):
        # AIが生成する列数がヘッダー数と違う可能性を考慮
        for c_idx, cell_text in enumerate(row_data):
            if c_idx < cols_count:
                cell = table.cell(r_idx + 1, c_idx)
                cell.text = sanitize_text(str(cell_text))
                # フォント設定
                p = cell.text_frame.paragraphs[0]
                p.font.name = "BIZ UDPGothic"
                p.font.size = Pt(12)

    # --- レイアウト調整（ここが重要） ---
    # 各行に固定の高さを設定し、テキストの折り返しを確実にする
    # 内容の量によらず、行の高さを固定することでレイアウトの破綻を防ぐ
    for row in table.rows:
        row.height = Inches(0.4) # 少し余裕のある高さを設定

    for row in table.rows:
        for cell in row.cells:
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE # 上下中央揃え
            if not cell.text_frame.word_wrap:
                cell.text_frame.word_wrap = True # 自動折り返しを有効化


# "static"ディレクトリ内のファイルを配信し、ルートパス("/")へのアクセスでindex.htmlを返す設定
app.mount("/", StaticFiles(directory="static", html=True), name="static")